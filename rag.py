"""
INVINCIBLE RAG core engine for ingesting study materials, storing chunked
content in ChromaDB, retrieving and reranking relevant context, managing
persistent conversation memory, and generating Gemini answers.

Key classes/methods:
- InvincibleRAG: main engine
- ingest_file(): parse and store supported files
- retrieve()/rerank(): hybrid retrieval and cross-encoder ranking
- generate_answer(): build prompt and call Gemini
- record_feedback(): store student corrections for future retrieval

How to run:
- Fill in `.env`
- Import `InvincibleRAG` from this file
- Use it directly or via `streamlit run app.py`
"""

from __future__ import annotations

import hashlib
import io
import os
import re
import sqlite3
import shutil
import time
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from uuid import uuid4

import chromadb
import fitz
import google.generativeai as genai
import pandas as pd
import pytesseract
from docx import Document
from dotenv import load_dotenv
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from PIL import Image
from pptx import Presentation
from sentence_transformers import CrossEncoder


load_dotenv()


def _utc_now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def _safe_int(value: Optional[str], default: int) -> int:
    try:
        return int(value) if value is not None else default
    except (TypeError, ValueError):
        return default


class InvincibleRAG:
    def __init__(self) -> None:
        self.gemini_api_key = os.getenv("GEMINI_API_KEY", "").strip()
        self.chroma_persist_dir = os.getenv("CHROMA_PERSIST_DIR", "./chroma_db")
        self.collection_name = os.getenv("COLLECTION_NAME", "student_rag")
        self.feedback_collection_name = os.getenv("FEEDBACK_COLLECTION", "rag_feedback")
        self.memory_collection_name = os.getenv("MEMORY_COLLECTION", "rag_memory")
        self.embedding_model_name = os.getenv("EMBEDDING_MODEL", "models/gemini-embedding-001")
        self.generation_model_name = os.getenv(
            "GENERATION_MODEL", "models/gemini-flash-lite-latest"
        )
        self.reranker_model_name = os.getenv(
            "RERANKER_MODEL", "cross-encoder/ms-marco-MiniLM-L-6-v2"
        )
        self.chunk_size = _safe_int(os.getenv("CHUNK_SIZE"), 800)
        self.chunk_overlap = _safe_int(os.getenv("CHUNK_OVERLAP"), 150)
        self.top_k_retrieval = _safe_int(os.getenv("TOP_K_RETRIEVAL"), 10)
        self.top_k_reranked = _safe_int(os.getenv("TOP_K_RERANKED"), 4)
        self.max_memory_turns = _safe_int(os.getenv("MAX_MEMORY_TURNS"), 10)
        self.ocr_min_text_threshold = _safe_int(os.getenv("OCR_MIN_TEXT_THRESHOLD"), 40)

        if not self.gemini_api_key:
            raise ValueError("GEMINI_API_KEY not found in environment.")

        genai.configure(api_key=self.gemini_api_key)

        self.client = self._initialize_chroma_client()
        self.collection, self.feedback_collection, self.memory_collection = self._initialize_collections()

        self.embedding_model_candidates = self._build_embedding_model_candidates(
            self.embedding_model_name
        )
        self.embeddings = self._initialize_embeddings()
        try:
            self.reranker = CrossEncoder(self.reranker_model_name)
        except (ImportError, OSError, RuntimeError):
            self.reranker = None
        self.generation_model_candidates = self._build_generation_model_candidates(
            self.generation_model_name
        )
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
            separators=["\n\n", "\n", ". ", "! ", "? ", " ", ""],
        )
        self.tesseract_available = shutil.which("tesseract") is not None

    def _initialize_chroma_client(self) -> chromadb.PersistentClient:
        try:
            return chromadb.PersistentClient(path=self.chroma_persist_dir)
        except sqlite3.OperationalError as exc:
            if "collections.topic" not in str(exc):
                raise

            self._reset_chroma_store()
            return chromadb.PersistentClient(path=self.chroma_persist_dir)

    def _initialize_collections(self) -> Tuple[Any, Any, Any]:
        try:
            return (
                self.client.get_or_create_collection(
                    name=self.collection_name,
                    metadata={"hnsw:space": "cosine"},
                ),
                self.client.get_or_create_collection(
                    name=self.feedback_collection_name,
                    metadata={"hnsw:space": "cosine"},
                ),
                self.client.get_or_create_collection(
                    name=self.memory_collection_name,
                    metadata={"hnsw:space": "cosine"},
                ),
            )
        except sqlite3.OperationalError as exc:
            if "collections.topic" not in str(exc):
                raise

            self._reset_chroma_store()
            self.client = chromadb.PersistentClient(path=self.chroma_persist_dir)
            return (
                self.client.get_or_create_collection(
                    name=self.collection_name,
                    metadata={"hnsw:space": "cosine"},
                ),
                self.client.get_or_create_collection(
                    name=self.feedback_collection_name,
                    metadata={"hnsw:space": "cosine"},
                ),
                self.client.get_or_create_collection(
                    name=self.memory_collection_name,
                    metadata={"hnsw:space": "cosine"},
                ),
            )

    def _reset_chroma_store(self) -> None:
        store_path = Path(self.chroma_persist_dir)
        if not store_path.exists():
            store_path.mkdir(parents=True, exist_ok=True)
            return

        backup_path = store_path.parent / (
            f"{store_path.name}.backup-{datetime.now(timezone.utc).strftime('%Y%m%d%H%M%S')}-{uuid4().hex[:8]}"
        )
        if backup_path.exists():
            shutil.rmtree(backup_path)
        shutil.move(str(store_path), str(backup_path))
        store_path.mkdir(parents=True, exist_ok=True)

    def _clean_text(self, text: str) -> str:
        text = text.replace("\x00", " ")
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    def _build_embedding_model_candidates(self, configured_model: str) -> List[str]:
        candidates: List[str] = []
        normalized = (configured_model or "").strip()
        if normalized:
            candidates.append(normalized)
            if normalized.startswith("models/"):
                candidates.append(normalized.replace("models/", "", 1))
        for fallback in [
            "models/gemini-embedding-001",
            "gemini-embedding-001",
            "models/gemini-embedding-2-preview",
            "gemini-embedding-2-preview",
        ]:
            if fallback not in candidates:
                candidates.append(fallback)
        return candidates

    def _initialize_embeddings(self) -> GoogleGenerativeAIEmbeddings:
        last_error: Optional[Exception] = None
        for model_name in self.embedding_model_candidates:
            try:
                embeddings = GoogleGenerativeAIEmbeddings(
                    model=model_name,
                    google_api_key=self.gemini_api_key,
                )
                test_vector = embeddings.embed_query("embedding connectivity check")
                if test_vector:
                    self.embedding_model_name = model_name
                    return embeddings
            except Exception as exc:
                last_error = exc
                continue
        raise RuntimeError(
            "Unable to initialize a Gemini embedding model. "
            f"Tried: {', '.join(self.embedding_model_candidates)}. Last error: {last_error}"
        ) from last_error

    def _build_generation_model_candidates(self, configured_model: str) -> List[str]:
        candidates: List[str] = []
        normalized = (configured_model or "").strip()
        if normalized:
            candidates.append(normalized)
            if normalized.startswith("models/"):
                candidates.append(normalized.replace("models/", "", 1))
        for fallback in [
            "models/gemini-flash-lite-latest",
            "gemini-flash-lite-latest",
            "models/gemini-2.5-flash-lite",
            "gemini-2.5-flash-lite",
            "models/gemma-3-1b-it",
            "gemma-3-1b-it",
        ]:
            if fallback not in candidates:
                candidates.append(fallback)
        return candidates

    def _batch_embed_documents(self, texts: List[str]) -> List[List[float]]:
        all_embeddings: List[List[float]] = []
        for start in range(0, len(texts), 50):
            batch = texts[start : start + 50]
            all_embeddings.extend(self.embeddings.embed_documents(batch))
            if start + 50 < len(texts):
                time.sleep(0.5)
        return all_embeddings

    def _embed_query(self, query: str) -> List[float]:
        return self.embeddings.embed_query(query)

    def _ocr_image_with_tesseract(self, image: Image.Image) -> str:
        if not self.tesseract_available:
            return ""
        try:
            return pytesseract.image_to_string(image).strip()
        except Exception:
            return ""

    def _ocr_image_with_gemini(self, image: Image.Image, context_label: str = "image") -> str:
        prompt = (
            f"Extract all readable text from this {context_label}. "
            "Preserve dates, headings, table rows, labels, and important formatting where possible. "
            "Return only the extracted text. If there is no readable text, return exactly: NO_TEXT"
        )
        last_error: Optional[Exception] = None
        rgb_image = image.convert("RGB")
        for model_name in self.generation_model_candidates:
            try:
                model = genai.GenerativeModel(model_name)
                response = model.generate_content(
                    [prompt, rgb_image],
                    generation_config=genai.types.GenerationConfig(
                        temperature=0.0,
                        top_p=0.8,
                        top_k=32,
                        max_output_tokens=2048,
                    ),
                )
                text = self._clean_text(getattr(response, "text", "") or "")
                if text.upper() == "NO_TEXT":
                    return ""
                if text:
                    return text
            except Exception as exc:
                last_error = exc
                continue
        if last_error:
            raise RuntimeError(f"Vision OCR failed for {context_label}: {last_error}") from last_error
        return ""

    def _ocr_image(self, image: Image.Image, context_label: str = "image") -> str:
        text = self._ocr_image_with_tesseract(image)
        if text:
            return text
        return self._ocr_image_with_gemini(image, context_label=context_label)

    def _render_pdf_page_image(self, page: fitz.Page, scale: float = 2.0) -> Image.Image:
        pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
        return Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")

    def _extract_pdf(self, file_bytes: bytes) -> List[Dict[str, Any]]:
        pages: List[Dict[str, Any]] = []
        with fitz.open(stream=file_bytes, filetype="pdf") as pdf:
            for page_index, page in enumerate(pdf, start=1):
                page_text = self._clean_text(page.get_text("text") or "")
                image_texts: List[str] = []
                for image_info in page.get_images(full=True):
                    xref = image_info[0]
                    try:
                        image_dict = pdf.extract_image(xref)
                        image_bytes = image_dict.get("image")
                        if not image_bytes:
                            continue
                        with Image.open(io.BytesIO(image_bytes)) as image:
                            ocr_text = self._ocr_image(image, context_label=f"PDF page {page_index} embedded image")
                        if ocr_text:
                            image_texts.append(ocr_text)
                    except Exception:
                        continue
                page_ocr_text = ""
                if len(page_text) < self.ocr_min_text_threshold and not image_texts:
                    try:
                        page_image = self._render_pdf_page_image(page)
                        page_ocr_text = self._ocr_image(page_image, context_label=f"PDF page {page_index}")
                    except Exception:
                        page_ocr_text = ""
                combined = f"PAGE {page_index}:\n{page_text}"
                if image_texts:
                    combined += "\n\nOCR TEXT FROM IMAGES:\n" + "\n".join(image_texts)
                if page_ocr_text:
                    combined += "\n\nOCR TEXT FROM PAGE RENDER:\n" + page_ocr_text
                pages.append({"label": page_index, "text": self._clean_text(combined)})
        return pages

    def _extract_pptx(self, file_bytes: bytes) -> List[Dict[str, Any]]:
        slides: List[Dict[str, Any]] = []
        presentation = Presentation(io.BytesIO(file_bytes))
        for slide_index, slide in enumerate(presentation.slides, start=1):
            chunks: List[str] = [f"SLIDE {slide_index}:"]
            for shape in slide.shapes:
                try:
                    if getattr(shape, "has_text_frame", False) and shape.text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            chunks.append(text)
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                            if row_text:
                                chunks.append(row_text)
                    if getattr(shape, "shape_type", None) == 13:
                        alt_text = getattr(shape, "alternative_text", "") or ""
                        if alt_text.strip():
                            chunks.append(f"IMAGE ALT TEXT: {alt_text.strip()}")
                except Exception:
                    continue
            slides.append({"label": slide_index, "text": self._clean_text("\n".join(chunks))})
        return slides

    def _extract_docx(self, file_bytes: bytes) -> List[Dict[str, Any]]:
        doc = Document(io.BytesIO(file_bytes))
        parts: List[str] = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style_name = getattr(paragraph.style, "name", "") if paragraph.style else ""
            if style_name.lower().startswith("heading"):
                parts.append(f"## HEADING: {text}")
            else:
                parts.append(text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)
        return [{"label": 1, "text": self._clean_text("\n".join(parts))}]

    def _extract_text_file(self, file_bytes: bytes) -> List[Dict[str, Any]]:
        return [{"label": 1, "text": self._clean_text(file_bytes.decode("utf-8", errors="ignore"))}]

    def _extract_csv(self, file_bytes: bytes) -> List[Dict[str, Any]]:
        df = pd.read_csv(io.BytesIO(file_bytes))
        description = (
            f"This dataset has {len(df)} rows and {len(df.columns)} columns: "
            f"{', '.join(map(str, df.columns.tolist()))}.\n"
            f"First 5 rows:\n{df.head().to_string(index=False)}"
        )
        return [{"label": 1, "text": self._clean_text(description)}]

    def _extract_image(self, file_bytes: bytes) -> List[Dict[str, Any]]:
        with Image.open(io.BytesIO(file_bytes)) as image:
            text = self._ocr_image(image, context_label="uploaded image")
        if not text:
            text = "Image with no readable text"
        return [{"label": 1, "text": self._clean_text(text)}]

    def _extract_content(self, extension: str, file_bytes: bytes) -> List[Dict[str, Any]]:
        extension = extension.lower()
        if extension == ".pdf":
            return self._extract_pdf(file_bytes)
        if extension == ".pptx":
            return self._extract_pptx(file_bytes)
        if extension == ".docx":
            return self._extract_docx(file_bytes)
        if extension in {".txt", ".md"}:
            return self._extract_text_file(file_bytes)
        if extension == ".csv":
            return self._extract_csv(file_bytes)
        if extension in {".png", ".jpg", ".jpeg", ".bmp", ".webp"}:
            return self._extract_image(file_bytes)
        raise ValueError(f"Unsupported file type: {extension}")

    def ingest_file(self, file_path: str, file_bytes: bytes) -> int:
        filename = os.path.basename(file_path)
        extension = os.path.splitext(filename)[1].lower()
        try:
            content_hash = hashlib.md5((filename.encode("utf-8") + file_bytes)).hexdigest()
            existing = self.collection.get(where={"document_id": content_hash}, include=["metadatas"])
            if existing.get("ids"):
                return 0

            extracted_sections = self._extract_content(extension, file_bytes)
            chunk_payloads: List[Tuple[str, Dict[str, Any]]] = []
            ingested_at = _utc_now_iso()

            for section in extracted_sections:
                section_label = section["label"]
                section_text = self._clean_text(section["text"])
                if not section_text:
                    continue
                section_chunks = self.text_splitter.split_text(section_text)
                total_chunks = len(section_chunks)
                for chunk_index, chunk_text in enumerate(section_chunks, start=1):
                    chunk_id = str(uuid4())
                    metadata = {
                        "source": filename,
                        "chunk_id": chunk_id,
                        "ingested_at": ingested_at,
                        "file_type": extension,
                        "chunk_index": chunk_index,
                        "total_chunks": total_chunks,
                        "page_or_slide": section_label,
                        "document_id": content_hash,
                    }
                    chunk_payloads.append((chunk_text, metadata))

            if not chunk_payloads:
                return 0

            texts = [chunk_text for chunk_text, _ in chunk_payloads]
            embeddings = self._batch_embed_documents(texts)

            for start in range(0, len(chunk_payloads), 50):
                batch_payloads = chunk_payloads[start : start + 50]
                batch_embeddings = embeddings[start : start + 50]
                self.collection.upsert(
                    ids=[metadata["chunk_id"] for _, metadata in batch_payloads],
                    documents=[chunk_text for chunk_text, _ in batch_payloads],
                    metadatas=[metadata for _, metadata in batch_payloads],
                    embeddings=batch_embeddings,
                )
                if start + 50 < len(chunk_payloads):
                    time.sleep(0.5)
            return len(chunk_payloads)
        except Exception as exc:
            raise RuntimeError(f"Failed to ingest {filename}: {exc}") from exc

    def _dense_retrieve(
        self, query: str, top_k: int
    ) -> List[Tuple[str, Dict[str, Any], float]]:
        query_embedding = self._embed_query(query)
        results = self.collection.query(
            query_embeddings=[query_embedding],
            n_results=top_k,
            include=["documents", "metadatas", "distances"],
        )
        documents = results.get("documents", [[]])[0]
        metadatas = results.get("metadatas", [[]])[0]
        distances = results.get("distances", [[]])[0]
        return list(zip(documents, metadatas, distances))

    def _keyword_retrieve(
        self, query: str, top_k: int
    ) -> List[Tuple[str, Dict[str, Any], float]]:
        keywords = [word.lower() for word in re.findall(r"\w+", query) if len(word) > 2]
        deduped: Dict[str, Tuple[str, Dict[str, Any], float]] = {}
        for keyword in keywords[:10]:
            try:
                results = self.collection.get(
                    where_document={"$contains": keyword},
                    include=["documents", "metadatas"],
                )
            except Exception:
                continue
            for doc_id, document, metadata in zip(
                results.get("ids", []),
                results.get("documents", []),
                results.get("metadatas", []),
            ):
                if doc_id not in deduped:
                    deduped[doc_id] = (document, metadata, 1.0)
                if len(deduped) >= top_k:
                    break
            if len(deduped) >= top_k:
                break
        return list(deduped.values())[:top_k]

    def _retrieve_feedback(self, query: str, top_k: int = 3) -> List[Tuple[str, Dict[str, Any], float]]:
        if self.feedback_collection.count() == 0:
            return []
        query_embedding = self._embed_query(query)
        results = self.feedback_collection.query(
            query_embeddings=[query_embedding],
            n_results=top_k,
            include=["documents", "metadatas", "distances"],
        )
        documents = results.get("documents", [[]])[0]
        metadatas = results.get("metadatas", [[]])[0]
        distances = results.get("distances", [[]])[0]
        return list(zip(documents, metadatas, distances))

    def retrieve(self, query: str, top_k: Optional[int] = None) -> List[Tuple[str, Dict[str, Any], float]]:
        max_results = top_k or self.top_k_retrieval
        dense_results = self._dense_retrieve(query, max_results)
        keyword_results = self._keyword_retrieve(query, max_results)
        feedback_results = self._retrieve_feedback(query, 3)

        merged: Dict[str, Tuple[str, Dict[str, Any], float]] = {}
        for document, metadata, distance in dense_results:
            merged[metadata["chunk_id"]] = (document, metadata, float(distance))
        for document, metadata, score in keyword_results:
            chunk_id = metadata["chunk_id"]
            if chunk_id not in merged:
                merged[chunk_id] = (document, metadata, float(score))

        merged_results = list(merged.values())[:max_results]
        feedback_prefix: List[Tuple[str, Dict[str, Any], float]] = []
        for document, metadata, distance in feedback_results:
            feedback_text = f"CORRECTION FROM STUDENT:\n{document}"
            feedback_meta = dict(metadata)
            feedback_meta.setdefault("source", "Student Feedback")
            feedback_meta.setdefault("page_or_slide", "N/A")
            feedback_meta["chunk_id"] = f"feedback-{metadata.get('query_hash', uuid4().hex)}"
            feedback_prefix.append((feedback_text, feedback_meta, float(distance)))

        return (feedback_prefix + merged_results)[:max_results]

    def rerank(
        self,
        query: str,
        candidates: List[Tuple[str, Dict[str, Any], float]],
        top_k: Optional[int] = None,
    ) -> List[Tuple[str, Dict[str, Any], float]]:
        if not candidates:
            return []
        rerank_k = top_k or self.top_k_reranked
        if self.reranker is None:
            return candidates[:rerank_k]
        pairs = [(query, candidate[0]) for candidate in candidates]
        scores = self.reranker.predict(pairs)
        reranked = [
            (document, metadata, float(score))
            for (document, metadata, _), score in zip(candidates, scores)
        ]
        reranked.sort(key=lambda item: item[2], reverse=True)
        return reranked[:rerank_k]

    def _get_next_turn_index(self, session_id: str) -> int:
        memory_records = self.memory_collection.get(
            where={"session_id": session_id},
            include=["metadatas"],
        )
        metadatas = memory_records.get("metadatas", [])
        if not metadatas:
            return 0
        return max(int(metadata.get("turn_index", 0)) for metadata in metadatas) + 1

    def save_memory_turn(self, session_id: str, role: str, content: str) -> None:
        turn_index = self._get_next_turn_index(session_id)
        record_id = str(uuid4())
        self.memory_collection.upsert(
            ids=[record_id],
            documents=[content],
            metadatas=[
                {
                    "session_id": session_id,
                    "role": role,
                    "timestamp": _utc_now_iso(),
                    "turn_index": turn_index,
                }
            ],
            embeddings=[self._embed_query(content[:8000])],
        )
        records = self.memory_collection.get(
            where={"session_id": session_id},
            include=["metadatas"],
        )
        entries = list(zip(records.get("ids", []), records.get("metadatas", [])))
        entries.sort(key=lambda item: int(item[1].get("turn_index", 0)))
        overflow = len(entries) - self.max_memory_turns
        if overflow > 0:
            self.memory_collection.delete(ids=[record_id for record_id, _ in entries[:overflow]])

    def get_memory(self, session_id: str) -> List[Dict[str, str]]:
        records = self.memory_collection.get(
            where={"session_id": session_id},
            include=["documents", "metadatas"],
        )
        combined = list(
            zip(records.get("documents", []), records.get("metadatas", []))
        )
        combined.sort(key=lambda item: int(item[1].get("turn_index", 0)))
        trimmed = combined[-self.max_memory_turns :]
        return [
            {"role": metadata.get("role", "user"), "content": document}
            for document, metadata in trimmed
        ]

    def format_memory_as_string(self, session_id: str) -> str:
        memory = self.get_memory(session_id)
        if not memory:
            return ""
        lines = []
        for turn in memory:
            speaker = "User" if turn["role"] == "user" else "Assistant"
            lines.append(f"{speaker}: {turn['content']}")
        return "\n".join(lines)

    def build_prompt(
        self,
        query: str,
        context_chunks: List[Tuple[str, Dict[str, Any], float]],
        memory_str: str,
        feedback_str: str,
    ) -> str:
        formatted_chunks = []
        for index, (chunk_text, metadata, score) in enumerate(context_chunks, start=1):
            source = metadata.get("source", "Unknown")
            page_or_slide = metadata.get("page_or_slide", "N/A")
            formatted_chunks.append(
                f"{index}. [Source: {source}, Page/Slide {page_or_slide}, Relevance: {score:.4f}]\n"
                f"{chunk_text}"
            )
        context_block = "\n\n".join(formatted_chunks) if formatted_chunks else "No retrieved context."
        return f"""
You are an expert AI study assistant helping a student understand
their study materials. You have access to the student's uploaded
documents and conversation history.

═══════════════════════════════════════════════════════
CONVERSATION HISTORY (most recent {self.max_memory_turns} turns):
═══════════════════════════════════════════════════════
{memory_str if memory_str else "No previous conversation."}

═══════════════════════════════════════════════════════
RETRIEVED DOCUMENT CONTEXT:
═══════════════════════════════════════════════════════
{context_block}

═══════════════════════════════════════════════════════
PREVIOUS CORRECTIONS / FEEDBACK:
═══════════════════════════════════════════════════════
{feedback_str if feedback_str else "None."}

═══════════════════════════════════════════════════════
STUDENT QUESTION:
═══════════════════════════════════════════════════════
{query}

═══════════════════════════════════════════════════════
INSTRUCTIONS:
═══════════════════════════════════════════════════════
1. Answer ONLY based on the provided context above.
2. If the answer is not in the context, say exactly:
   "I could not find this in your uploaded documents. Here is
    what I know generally: ..." — then give a general answer.
3. Always cite your sources: after each key claim write
   [Source: {{filename}}, Page/Slide {{n}}] inline.
4. Structure your answer clearly with:
     • A direct answer in the first sentence
     • Supporting details with citations
     • A concise summary at the end if the answer is long
5. If the question involves a concept, also explain it step
   by step in simple language.
6. Never fabricate information. Never say you cannot help.
7. Match the student's language (if they write in Hindi or
   another language, respond in that language too).
8. Answer the exact question first and stay tightly relevant.
9. Do not include unrelated implementation details, file paths,
   dataset structure, counts, or extra background unless the
   student explicitly asks for them.
10. For fact questions like "which dataset", "what model",
    "what algorithm", or "what is used", respond in 2-4 sentences
    maximum unless the student asks for a detailed explanation.
11. If multiple relevant details exist, present only the most
    important ones first, then add brief supporting context.
12. Never expose local file system paths unless the student
    explicitly asks for the path or location.
""".strip()

    def generate_answer(self, query: str, session_id: str) -> Dict[str, Any]:
        memory_str = self.format_memory_as_string(session_id)
        candidates = self.retrieve(query)
        top_chunks = self.rerank(query, candidates)
        feedback_candidates = self._retrieve_feedback(query, 3)
        feedback_str = "\n\n".join(
            f"CORRECTION FROM STUDENT:\n{document}" for document, _, _ in feedback_candidates
        )
        prompt = self.build_prompt(query, top_chunks, memory_str, feedback_str)
        response = None
        answer_model = self.generation_model_name
        last_error: Optional[Exception] = None
        for model_name in self.generation_model_candidates:
            try:
                model = genai.GenerativeModel(model_name)
                response = model.generate_content(
                    prompt,
                    generation_config=genai.types.GenerationConfig(
                        temperature=0.2,
                        top_p=0.85,
                        top_k=40,
                        max_output_tokens=2048,
                    ),
                )
                answer_model = model_name
                break
            except Exception as exc:
                last_error = exc
                continue
        if response is None:
            raise RuntimeError(
                "Unable to generate an answer with the available Gemini models. "
                f"Tried: {', '.join(self.generation_model_candidates)}. Last error: {last_error}"
            ) from last_error
        answer_text = getattr(response, "text", "") or "No answer generated."
        self.save_memory_turn(session_id, "user", query)
        self.save_memory_turn(session_id, "assistant", answer_text)
        return {
            "answer": answer_text,
            "sources": [metadata for _, metadata, _ in top_chunks],
            "chunks_used": [chunk_text for chunk_text, _, _ in top_chunks],
            "model": answer_model,
            "scores": [score for _, _, score in top_chunks],
        }

    def record_feedback(
        self,
        query: str,
        bad_answer: str,
        correction: str,
        session_id: str,
    ) -> None:
        feedback_text = (
            f"QUESTION: {query}\n"
            f"WRONG ANSWER: {bad_answer}\n"
            f"CORRECT ANSWER: {correction}"
        )
        feedback_id = str(uuid4())
        query_hash = hashlib.md5(query.encode("utf-8")).hexdigest()
        self.feedback_collection.upsert(
            ids=[feedback_id],
            documents=[feedback_text],
            metadatas=[
                {
                    "session_id": session_id,
                    "timestamp": _utc_now_iso(),
                    "query_hash": query_hash,
                }
            ],
            embeddings=[self._embed_query(feedback_text)],
        )

    def list_ingested_files(self) -> List[str]:
        records = self.collection.get(include=["metadatas"])
        filenames = {
            metadata.get("source", "")
            for metadata in records.get("metadatas", [])
            if metadata.get("source")
        }
        return sorted(filenames)

    def delete_file(self, filename: str) -> None:
        records = self.collection.get(where={"source": filename})
        ids = records.get("ids", [])
        if ids:
            self.collection.delete(ids=ids)

    def get_stats(self) -> Dict[str, int]:
        records = self.collection.get(include=["metadatas"])
        total_chunks = self.collection.count()
        unique_files = {
            metadata.get("source", "")
            for metadata in records.get("metadatas", [])
            if metadata.get("source")
        }
        memory_records = self.memory_collection.get(include=["metadatas"])
        sessions = {
            metadata.get("session_id", "")
            for metadata in memory_records.get("metadatas", [])
            if metadata.get("session_id")
        }
        return {
            "total_chunks": total_chunks,
            "total_files": len(unique_files),
            "total_feedback": self.feedback_collection.count(),
            "total_sessions": len(sessions),
        }
